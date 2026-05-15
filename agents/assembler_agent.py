"""
Agent 6 — Deck Assembler Agent

Takes DeckContent + ArchitectureDescription and assembles the 46-slide
PPTX following the required structure from the UC1 brief.
"""
from __future__ import annotations
import pathlib
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from models.schemas import DeckContent, ArchitectureDescription, FeatureServiceMatrixRow
import config
import re


def _strip_markdown(text: str) -> str:
    """Remove markdown syntax that leaks from LLM outputs into slide text."""
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)  # headings
    text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)        # bold/italic
    text = re.sub(r'`([^`]+)`', r'\1', text)                     # inline code
    text = re.sub(r'^\s*[-*]\s+\[', '• [', text, flags=re.MULTILINE)  # md list items
    text = re.sub(r'\n{3,}', '\n\n', text)                       # excess blank lines
    return text.strip()


# ── Brand colours ─────────────────────────────────────────────────────────────
PURPLE  = RGBColor(0x5B, 0x2D, 0x8E)
LPURPLE = RGBColor(0xD4, 0xA8, 0xF5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT   = RGBColor(0xF5, 0xF0, 0xFF)
GREY    = RGBColor(0x88, 0x88, 0x88)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_layout(prs: Presentation):
    return prs.slide_layouts[6]  # blank


def _add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def _set_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide builders ─────────────────────────────────────────────────────────────

def _add_cover(prs: Presentation, content: DeckContent) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, DARK)

    # Purple accent bar left
    _add_shape(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill_color=PURPLE)

    # Client/project label
    _add_textbox(slide, f"{content.client_name}  |  {content.project_name}",
                 Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 font_size=13, color=LPURPLE, bold=True)

    # Main title
    _add_textbox(slide, "AI-Powered RFP Response",
                 Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
                 font_size=36, bold=True, color=WHITE)

    # Subtitle
    _add_textbox(slide, f"Proposal — {content.project_name}",
                 Inches(0.8), Inches(3.3), Inches(10), Inches(0.5),
                 font_size=16, color=LPURPLE)

    # Meta info
    meta = f"Version: {content.version}   |   Date: {content.date}   |   Status: {content.status}"
    _add_textbox(slide, meta, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
                 font_size=11, color=GREY)

    _add_textbox(slide, f"Prepared by: {', '.join(content.authors)}",
                 Inches(0.8), Inches(6.0), Inches(10), Inches(0.4),
                 font_size=11, color=GREY)


def _add_divider(prs: Presentation, label: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, PURPLE)
    _add_shape(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.08), fill_color=WHITE)
    _add_textbox(slide, label,
                 Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                     font_size=16, color=LPURPLE, align=PP_ALIGN.CENTER)


def _add_content_slide(prs: Presentation, title: str, body: str,
                        bullets: List[str] | None = None) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, WHITE)
    # Purple header bar
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
    _add_textbox(slide, title, Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                 font_size=20, bold=True, color=WHITE)
    # Accent line
    _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)
    elif body:
        _add_textbox(slide, body, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
                     font_size=13, color=DARK)


def _add_table_slide(prs: Presentation, title: str,
                     headers: List[str], rows: List[List[str]],
                     page_label: str = "") -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, WHITE)

    # Header bar
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
    _add_textbox(slide, title + (f"  {page_label}" if page_label else ""),
                 Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                 font_size=20, bold=True, color=WHITE)
    _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

    if not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    table_left = Inches(0.3)
    table_top  = Inches(1.25)
    table_w    = Inches(12.7)
    table_h    = Inches(5.9)

    tbl = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h).table

    # Column widths (distribute evenly)
    col_w = int(table_w / n_cols)
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    # Header row
    for col_idx, h in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.font.bold  = True
        p.font.color.rgb = WHITE
        p.font.size  = Pt(11)
        p.alignment  = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows, start=1):
        fill = LIGHT if row_idx % 2 == 0 else WHITE
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx, col_idx)
            cell.text = str(val) if val is not None else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK


def _add_service_matrix(prs: Presentation,
                         feature_matrix: List[FeatureServiceMatrixRow],
                         all_service_names: List[str]) -> None:
    """Slide 5 — Feature × Service matrix."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
    _add_textbox(slide, "Feature & Service Matrix",
                 Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                 font_size=20, bold=True, color=WHITE)
    _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

    # Trim service names to fit
    svc_names = all_service_names[:8]  # max 8 columns to fit slide
    headers = ["#", "Feature Group"] + [s[:18] for s in svc_names]
    rows = []
    for row in feature_matrix:
        row_data = [str(row.item_num), row.feature[:35]]
        for svc in svc_names:
            row_data.append("X" if svc in row.services else "")
        rows.append(row_data)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  Inches(0.2), Inches(1.25),
                                  Inches(12.9), Inches(5.9)).table

    # Col widths
    tbl.columns[0].width = Inches(0.5)
    tbl.columns[1].width = Inches(3.0)
    rest = int((Inches(12.9) - Inches(3.5)) / max(len(svc_names), 1))
    for i in range(2, n_cols):
        tbl.columns[i].width = rest

    for col_idx, h in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows, start=1):
        fill = LIGHT if row_idx % 2 == 0 else WHITE
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = PURPLE if val == "X" else DARK
            p.font.bold = (val == "X")
            p.alignment = PP_ALIGN.CENTER if val == "X" else PP_ALIGN.LEFT


def _add_architecture_diagram(prs: Presentation, arch: ArchitectureDescription) -> None:
    """Slide 8 — High-Level Architecture (text-based diagram)."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
    _add_textbox(slide, "High-Level Architecture",
                 Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                 font_size=20, bold=True, color=WHITE)
    _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

    # Layers as boxes
    layers = [
        ("Clients / End Users", "HR Employees · Managers · Admins", Inches(0.3), Inches(1.3)),
        ("Internet Boundary / SSO", "Identity Provider (IdP) · MFA · OAuth 2.0", Inches(0.3), Inches(2.2)),
        (f"Platform Layer — {arch.platform_name}",
         "  ·  ".join(s.name[:22] for s in arch.services[:5]),
         Inches(0.3), Inches(3.1)),
        ("On-Premise Systems", arch.cloud_boundary, Inches(0.3), Inches(4.5)),
        ("Monitoring", "DynaTrace · Alert Notification · Automation Pilot", Inches(0.3), Inches(5.5)),
    ]

    colors = [DARK, PURPLE, RGBColor(0x4A, 0x0, 0x80), RGBColor(0x1A, 0x4D, 0x80), RGBColor(0x0, 0x5C, 0x35)]

    for (title, desc, left, top), color in zip(layers, colors):
        box = _add_shape(slide, left, top, Inches(12.7), Inches(0.75), fill_color=color)
        _add_textbox(slide, title, left + Inches(0.1), top + Inches(0.05),
                     Inches(4), Inches(0.35), font_size=11, bold=True, color=WHITE)
        _add_textbox(slide, desc, left + Inches(4.5), top + Inches(0.05),
                     Inches(8), Inches(0.6), font_size=10, color=WHITE)


def _add_plan_slide(prs: Presentation, content: DeckContent) -> None:
    """Slide 21 — High-level plan (Gantt-style)."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_background(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
    _add_textbox(slide, "High-Level Project Plan",
                 Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                 font_size=20, bold=True, color=WHITE)
    _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

    _add_textbox(slide, content.plan_note,
                 Inches(0.5), Inches(1.2), Inches(12), Inches(0.35),
                 font_size=10, color=GREY, italic=True)

    # Month header row
    total_months = max(p.start_month + p.duration_months for p in content.plan_phases) - 1
    month_w = Inches(10.5) / total_months
    label_col_w = Inches(2.5)
    gantt_left = Inches(0.5) + label_col_w

    for m in range(total_months):
        _add_textbox(slide, f"M{m+1}",
                     gantt_left + Emu(int(m * month_w)), Inches(1.65),
                     Emu(int(month_w)), Inches(0.3),
                     font_size=8, color=GREY, align=PP_ALIGN.CENTER)

    # Phase rows
    bar_colors = [PURPLE, RGBColor(0x4A, 0x90, 0xD9), RGBColor(0x40, 0xB0, 0x70),
                  RGBColor(0xE0, 0xB0, 0x40), RGBColor(0xD0, 0x40, 0x80), DARK]

    for i, phase in enumerate(content.plan_phases):
        row_top = Inches(2.0) + Emu(int(i * Inches(0.72)))
        _add_textbox(slide, phase.phase,
                     Inches(0.5), row_top, label_col_w, Inches(0.55),
                     font_size=10, color=DARK, bold=True)
        bar_left  = gantt_left + Emu(int((phase.start_month - 1) * month_w))
        bar_width = Emu(int(phase.duration_months * month_w))
        color = bar_colors[i % len(bar_colors)]
        _add_shape(slide, bar_left, row_top + Emu(int(Inches(0.08))),
                   bar_width, Inches(0.38), fill_color=color)


# ── Feature-solution table slides ─────────────────────────────────────────────

def _add_feature_solution_slides(prs: Presentation,
                                  feature_matrix: List[FeatureServiceMatrixRow]) -> None:
    """Slides 3-4: Feature-Solution tables, up to 8 rows per slide."""
    ROWS_PER_SLIDE = 8
    chunks = [feature_matrix[i:i+ROWS_PER_SLIDE]
              for i in range(0, len(feature_matrix), ROWS_PER_SLIDE)]

    for chunk_idx, chunk in enumerate(chunks):
        page_label = f"({chunk_idx + 1}/{len(chunks)})" if len(chunks) > 1 else ""
        headers = ["#", "Feature", "Req. IDs", "Proposed Solution"]
        rows = [
            [str(row.item_num),
             row.feature[:40],
             ", ".join(row.req_ids[:6]) + ("…" if len(row.req_ids) > 6 else ""),
             row.solution[:60]]
            for row in chunk
        ]
        _add_table_slide(prs, "Features — Proposed Solution", headers, rows, page_label)


# ── Service description slides ────────────────────────────────────────────────

def _add_service_description_slides(prs: Presentation,
                                     arch: ArchitectureDescription) -> None:
    """Slides 10-13: 3-4 services per slide."""
    SERVICES_PER_SLIDE = 4
    services = arch.services
    chunks = [services[i:i+SERVICES_PER_SLIDE]
              for i in range(0, len(services), SERVICES_PER_SLIDE)]

    for chunk_idx, chunk in enumerate(chunks):
        page_label = f"({chunk_idx + 1}/{len(chunks)})" if len(chunks) > 1 else ""
        slide = prs.slides.add_slide(_blank_layout(prs))
        _set_background(slide, WHITE)
        _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
        _add_textbox(slide,
                     f"Platform Services — Description  {page_label}",
                     Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                     font_size=20, bold=True, color=WHITE)
        _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)

        row_h = Inches(5.8) / max(len(chunk), 1)
        for svc_idx, svc in enumerate(chunk):
            top = Inches(1.25) + Emu(int(svc_idx * row_h))
            _add_shape(slide, Inches(0.2), top, Inches(12.9), Emu(int(row_h - Inches(0.1))),
                       fill_color=LIGHT)
            _add_textbox(slide, svc.name,
                         Inches(0.4), top + Inches(0.05), Inches(12.5), Inches(0.35),
                         font_size=13, bold=True, color=PURPLE)
            _add_textbox(slide, svc.description,
                         Inches(0.4), top + Inches(0.4), Inches(12.5),
                         Emu(int(row_h - Inches(0.55))),
                         font_size=11, color=DARK)


# ── Roles slides ──────────────────────────────────────────────────────────────

def _add_roles_slides(prs: Presentation, content: DeckContent) -> None:
    """Slides 17-19: Roles & Responsibilities tables."""
    ROLES_PER_SLIDE = 2
    roles = content.roles
    chunks = [roles[i:i+ROLES_PER_SLIDE] for i in range(0, len(roles), ROLES_PER_SLIDE)]

    for chunk_idx, chunk in enumerate(chunks):
        page_label = f"({chunk_idx + 1}/{len(chunks)})"
        headers = ["Role", "Responsibilities", "Outcomes", "Responsible Party"]
        rows = [
            [r.role,
             "\n".join(f"• {x}" for x in r.responsibilities),
             "\n".join(f"• {x}" for x in r.outcomes),
             r.responsible_party]
            for r in chunk
        ]
        _add_table_slide(prs, "Roles & Responsibilities", headers, rows, page_label)


# ── Assumptions slides ────────────────────────────────────────────────────────

def _add_assumptions_slides(prs: Presentation, content: DeckContent) -> None:
    """Slides 28-31: Assumptions, ~3 per slide."""
    from itertools import groupby
    by_cat = {}
    for a in content.assumptions:
        by_cat.setdefault(a.category, []).append(a.text)

    # Split into 4 slides
    items = [(cat, text) for cat, texts in by_cat.items() for text in texts]
    PER_SLIDE = max(1, len(items) // 4 + 1)
    chunks = [items[i:i+PER_SLIDE] for i in range(0, len(items), PER_SLIDE)]

    for chunk_idx, chunk in enumerate(chunks[:4]):
        page_label = f"({chunk_idx + 1}/4)"
        bullets = [f"[{cat}] {text}" for cat, text in chunk]
        _add_content_slide(prs, f"Assumptions  {page_label}", "", bullets=bullets)


# ── Credentials slides ────────────────────────────────────────────────────────

def _add_credentials_slides(prs: Presentation, content: DeckContent) -> None:
    """Slides 33-39: One per credential."""
    for cred in content.credentials[:7]:
        slide = prs.slides.add_slide(_blank_layout(prs))
        _set_background(slide, WHITE)
        _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=PURPLE)
        _add_textbox(slide, f"Credentials — {cred.industry}",
                     Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.7),
                     font_size=20, bold=True, color=WHITE)
        _add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=LPURPLE)
        _add_textbox(slide, cred.project,
                     Inches(0.5), Inches(1.25), Inches(12), Inches(0.5),
                     font_size=16, bold=True, color=PURPLE)
        _add_textbox(slide, cred.description,
                     Inches(0.5), Inches(1.9), Inches(12), Inches(1.2),
                     font_size=12, color=DARK)
        bullets = [f"• {o}" for o in cred.outcomes]
        _add_textbox(slide, "\n".join(bullets),
                     Inches(0.5), Inches(3.2), Inches(12), Inches(3.5),
                     font_size=12, color=DARK)


# ── Appendix slides ───────────────────────────────────────────────────────────

def _add_appendix_slides(prs: Presentation, content: DeckContent) -> None:
    """Slides 41-46: Appendix items."""
    bullets_per_slide = 3
    items = content.appendix_items
    chunks = [items[i:i+bullets_per_slide] for i in range(0, len(items), bullets_per_slide)]
    for chunk_idx, chunk in enumerate(chunks[:6]):
        page_label = f"({chunk_idx + 1}/{len(chunks)})"
        _add_content_slide(prs, f"Appendix  {page_label}", "", bullets=chunk)


# ── Public interface ──────────────────────────────────────────────────────────

def run(content: DeckContent, arch: ArchitectureDescription) -> pathlib.Path:
    """Assemble the 46-slide PPTX and return the output path."""
    print("  [Assembler] Building PPTX …")
    prs = _new_prs()

    all_service_names = [s.name for s in arch.services]

    # ── Slides 1: Cover ──────────────────────────────────────────────────────
    print("  [Assembler] Slide 1: Cover")
    _add_cover(prs, content)

    # ── Slide 2: Scope divider ────────────────────────────────────────────────
    print("  [Assembler] Slide 2: Scope divider")
    _add_divider(prs, "Scope", "Our Understanding of the RFP")

    # ── Slides 3-4: Feature-solution tables ───────────────────────────────────
    print("  [Assembler] Slides 3-4: Feature-solution tables")
    _add_feature_solution_slides(prs, arch.feature_matrix)

    # ── Slide 5: Feature & Service Matrix ────────────────────────────────────
    print("  [Assembler] Slide 5: Service matrix")
    _add_service_matrix(prs, arch.feature_matrix, all_service_names)

    # ── Slide 6: Principles alignment ────────────────────────────────────────
    print("  [Assembler] Slide 6: Principles alignment")
    _add_content_slide(prs, "Principles Alignment", content.principles_text)

    # ── Slide 7: Architecture divider ────────────────────────────────────────
    print("  [Assembler] Slide 7: Architecture divider")
    _add_divider(prs, "Architecture")

    # ── Slide 8: High-level architecture diagram ──────────────────────────────
    print("  [Assembler] Slide 8: Architecture diagram")
    _add_architecture_diagram(prs, arch)

    # ── Slide 9: Architecture services list ──────────────────────────────────
    print("  [Assembler] Slide 9: Services list")
    bullets = [f"{s.name}  ({s.platform})" for s in arch.services]
    _add_content_slide(prs, f"Platform Services — {arch.platform_name}", arch.narrative, bullets=bullets)

    # ── Slides 10-13: Service descriptions ───────────────────────────────────
    print("  [Assembler] Slides 10-13: Service descriptions")
    _add_service_description_slides(prs, arch)

    # ── Slide 14: Approach divider ────────────────────────────────────────────
    print("  [Assembler] Slide 14: Approach divider")
    _add_divider(prs, "Approach")

    # ── Slide 15: Project setup & team ───────────────────────────────────────
    print("  [Assembler] Slide 15: Project setup")
    _add_content_slide(prs, "Project Setup & Team",
                       _strip_markdown(f"Delivery: {content.delivery_method}\n\nLocations: {content.locations}\n\n{content.team_structure}"))

    # ── Slide 16: Governance ──────────────────────────────────────────────────
    print("  [Assembler] Slide 16: Governance")
    _add_content_slide(prs, "Governance & Communication", _strip_markdown(content.governance))

    # ── Slides 17-19: Roles & Responsibilities ────────────────────────────────
    print("  [Assembler] Slides 17-19: R&R")
    _add_roles_slides(prs, content)

    # ── Slide 20: Plan divider ────────────────────────────────────────────────
    print("  [Assembler] Slide 20: Plan divider")
    _add_divider(prs, "Plan")

    # ── Slide 21: High-level plan ─────────────────────────────────────────────
    print("  [Assembler] Slide 21: Plan timeline")
    _add_plan_slide(prs, content)

    # ── Slide 22: Estimates divider (optional) ────────────────────────────────
    print("  [Assembler] Slide 22: Estimates divider")
    _add_divider(prs, "Estimates", "(Preliminary — to be finalised after Discovery)")

    # ── Slides 23-26: Placeholder estimates ──────────────────────────────────
    for title in [
        "Features Breakdown by Artefacts",
        "High-Level Effort — Implementation",
        "High-Level Effort — Hypercare",
        "High-Level Effort — Total",
    ]:
        _add_content_slide(prs, title,
                           "Detailed estimates to be provided after Discovery & Planning phase.",
                           bullets=["Effort sizing will be based on finalised requirements",
                                    "RICEFW object counts to be confirmed",
                                    "Resource allocation subject to project kick-off alignment"])

    # ── Slide 27: Assumptions divider ─────────────────────────────────────────
    print("  [Assembler] Slides 27-31: Assumptions")
    _add_divider(prs, "Assumptions")

    # ── Slides 28-31: Assumptions ─────────────────────────────────────────────
    _add_assumptions_slides(prs, content)

    # ── Slide 32: Credentials divider ─────────────────────────────────────────
    print("  [Assembler] Slides 32-39: Credentials")
    _add_divider(prs, "Credentials")

    # ── Slides 33-39: Credentials ─────────────────────────────────────────────
    _add_credentials_slides(prs, content)

    # ── Slide 40: Appendix divider ────────────────────────────────────────────
    print("  [Assembler] Slides 40-46: Appendix")
    _add_divider(prs, "Appendix")

    # ── Slides 41-46: Appendix ───────────────────────────────────────────────
    _add_appendix_slides(prs, content)

    # ── Save ──────────────────────────────────────────────────────────────────
    output_path = config.OUTPUTS_DIR / "Response_Deck.pptx"
    try:
        prs.save(str(output_path))
    except PermissionError:
        from datetime import datetime
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = config.OUTPUTS_DIR / f"Response_Deck_{ts}.pptx"
        prs.save(str(output_path))
    print(f"  [Assembler] PPTX saved → {output_path.name}  ({len(prs.slides)} slides)")
    return output_path
