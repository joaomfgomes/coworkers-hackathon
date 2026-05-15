"""
Agent 1 — Document Ingestion Agent

Reads DOCX, PDF, or PPTX files and returns a structured IngestedDocument
with one DocumentSection per logical section / heading.

No LLM call needed here — pure document parsing.
"""
from __future__ import annotations
import pathlib
from typing import List
from models.schemas import DocumentSection, IngestedDocument


# ── DOCX parser ───────────────────────────────────────────────────────────────

def _parse_docx(path: pathlib.Path) -> List[DocumentSection]:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(path))
    sections: List[DocumentSection] = []
    current_title = "Introduction"
    current_body: List[str] = []
    section_idx = 0

    def _flush():
        nonlocal section_idx
        if current_body:
            sections.append(DocumentSection(
                source_file=path.name,
                section_id=str(section_idx),
                title=current_title,
                content="\n".join(current_body).strip(),
            ))
            section_idx += 1

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            _flush()
            current_title = text
            current_body = []
        else:
            current_body.append(text)

    _flush()

    # Also extract tables
    for tbl_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            sections.append(DocumentSection(
                source_file=path.name,
                section_id=f"table-{tbl_idx}",
                title=f"Table {tbl_idx + 1}",
                content="\n".join(" | ".join(r) for r in rows),
                tables=rows,
            ))

    return sections


# ── PDF parser ────────────────────────────────────────────────────────────────

def _parse_pdf(path: pathlib.Path) -> List[DocumentSection]:
    import pdfplumber

    sections: List[DocumentSection] = []
    current_title = "Section 1"
    current_body: List[str] = []
    section_idx = 0

    def _flush():
        nonlocal section_idx
        if current_body:
            sections.append(DocumentSection(
                source_file=path.name,
                section_id=str(section_idx),
                title=current_title,
                content="\n".join(current_body).strip(),
            ))
            section_idx += 1

    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            for line in text.split("\n"):
                stripped = line.strip()
                if not stripped:
                    continue
                # Heuristic: all-caps short lines or numbered headings are section titles
                if (len(stripped) < 80 and (
                    stripped.isupper() or
                    (stripped[0].isdigit() and "." in stripped[:4])
                )):
                    _flush()
                    current_title = stripped
                    current_body = []
                else:
                    current_body.append(stripped)

    _flush()
    return sections


# ── PPTX parser ───────────────────────────────────────────────────────────────

def _parse_pptx(path: pathlib.Path) -> List[DocumentSection]:
    from pptx import Presentation

    prs = Presentation(str(path))
    sections: List[DocumentSection] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = f"Slide {slide_num}"
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    if shape.shape_type == 13:   # title placeholder
                        title = t
                    else:
                        texts.append(t)
        if texts:
            sections.append(DocumentSection(
                source_file=path.name,
                section_id=f"slide-{slide_num}",
                title=title,
                content="\n".join(texts),
            ))

    return sections


# ── Public interface ──────────────────────────────────────────────────────────

def run(file_paths: List[str | pathlib.Path]) -> IngestedDocument:
    """
    Parse one or more RFP documents and return an IngestedDocument.

    Supported formats: .docx, .pdf, .pptx
    """
    all_sections: List[DocumentSection] = []

    for fp in file_paths:
        path = pathlib.Path(fp)
        suffix = path.suffix.lower()
        print(f"  [Ingestion] Parsing {path.name} ({suffix}) …")

        if suffix == ".docx":
            sections = _parse_docx(path)
        elif suffix == ".pdf":
            sections = _parse_pdf(path)
        elif suffix in (".pptx", ".ppt"):
            sections = _parse_pptx(path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")

        all_sections.extend(sections)
        print(f"  [Ingestion] → {len(sections)} sections extracted from {path.name}")

    return IngestedDocument(sections=all_sections)
